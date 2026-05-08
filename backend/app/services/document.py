import logging
import uuid
import os
import io
import fitz  # PyMuPDF
from typing import List
import google.generativeai as genai
from pptx import Presentation
from docx import Document as DocxDocument

from fastapi import HTTPException

from app.core.config import MAX_FILE_SIZE, GEMINI_API_KEY
from app.core.supabase import supabase

logger = logging.getLogger(__name__)

# Gemini API 설정
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# 사용할 모델 정의
CHAT_MODEL = "gemini-2.5-flash-lite" # Gemini 2.0 Flash Lite 정식 프리뷰 명칭
EMBEDDING_MODEL = "gemini-embedding-001"

MAGIC_NUMBERS = {
    ".pdf":  b"%PDF",
    ".pptx": b"PK\x03\x04",
    ".docx": b"PK\x03\x04",
    ".ppt":  b"\xd0\xcf\x11\xe0",
}

CONTENT_TYPES = {
    ".pdf":  "application/pdf",
    ".ppt":  "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".txt":  "text/plain; charset=utf-8",
}

ALLOWED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".docx", ".txt"}


def validate_file(original_name: str, contents: bytes) -> str:
    ext = os.path.splitext(original_name)[1].lower()
    if not original_name or ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="허용되지 않는 파일 형식입니다.")

    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="파일 크기는 50MB를 초과할 수 없습니다.")

    magic = MAGIC_NUMBERS.get(ext)
    if magic and not contents.startswith(magic):
        raise HTTPException(status_code=400, detail="파일 형식이 올바르지 않습니다.")

    if ext == ".txt":
        try:
            contents.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="파일 형식이 올바르지 않습니다.")

    return ext


def upload_document(original_name: str, contents: bytes, ext: str, student_id: str = None, category_id: int = None) -> int:
    file_path = f"{uuid.uuid4()}{ext}"

    try:
        content_type = CONTENT_TYPES[ext]
        supabase.storage.from_("documents").upload(file_path, contents, {"content-type": content_type})
    except Exception as e:
        logger.error("스토리지 업로드 실패: %s", e)
        raise HTTPException(status_code=500, detail="파일 저장 중 오류가 발생했습니다.")

    try:
        res = supabase.table("documents").insert({
            "student_id": student_id,
            "category_id": category_id,
            "original_file_name": original_name,
            "file_path": file_path,
            "parsing_status": "PENDING",
        }).execute()
        
        if not res.data:
            raise Exception("DB insert result is empty")
            
        return res.data[0]["id"]
    except Exception as e:
        logger.error("DB 저장 실패: %s", e)
        try:
            supabase.storage.from_("documents").remove([file_path])
        except Exception as cleanup_err:
            logger.error("스토리지 정리 실패: %s", cleanup_err)
        raise HTTPException(status_code=500, detail="파일 정보 저장 중 오류가 발생했습니다.")


async def process_document_rag(document_id: int):
    """
    RAG 인제스션 파이프라인 (Gemini 버전):
    1. 파일 다운로드
    2. 텍스트 추출
    3. 청킹
    4. Gemini 임베딩 생성
    5. 벡터 DB 저장
    6. 상태 업데이트
    """
    try:
        # 1. 문서 정보 가져오기
        res = supabase.table("documents").select("*").eq("id", document_id).single().execute()
        doc_data = res.data
        if not doc_data:
            logger.error(f"문서 ID {document_id}를 찾을 수 없습니다.")
            return

        file_path = doc_data["file_path"]
        ext = os.path.splitext(file_path)[1].lower()

        # 2. 파일 다운로드
        file_bytes = supabase.storage.from_("documents").download(file_path)

        # 3. 텍스트 추출
        text = ""
        if ext == ".pdf":
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
        elif ext == ".txt":
            text = file_bytes.decode("utf-8")
        elif ext in (".pptx",):
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text += para.text + "\n"
        elif ext in (".docx",):
            docx = DocxDocument(io.BytesIO(file_bytes))
            for para in docx.paragraphs:
                text += para.text + "\n"
        else:
            logger.warning(f"{ext} 형식은 텍스트 추출을 지원하지 않습니다.")
            text = f"[{doc_data['original_file_name']} - 텍스트 추출 미지원 형식]"

        if not text.strip():
            raise Exception("추출된 텍스트가 없습니다.")

        # 4. 청킹
        chunk_size = 1000  # Gemini는 컨텍스트 윈도우가 크므로 약간 더 크게 잡아도 좋습니다.
        chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

        # 5. 임베딩 생성 및 저장
        for i, chunk_content in enumerate(chunks):
            # Gemini 임베딩 생성 (models/embedding-001 사용)
            embedding_res = genai.embed_content(
                model=EMBEDDING_MODEL,
                content=chunk_content,
                task_type="retrieval_document"
            )
            embedding_vector = embedding_res['embedding']

            # document_chunks 테이블에 저장
            supabase.table("document_chunks").insert({
                "document_id": document_id,
                "content": chunk_content,
                "embedding": embedding_vector,
                "chunk_index": i
            }).execute()

        # 6. 상태 업데이트
        supabase.table("documents").update({"parsing_status": "COMPLETED"}).eq("id", document_id).execute()
        logger.info(f"문서 {document_id} RAG 처리 완료 (Gemini 사용)")

    except Exception as e:
        logger.error(f"문서 {document_id} RAG 처리 중 오류: {e}")
        supabase.table("documents").update({"parsing_status": "FAILED"}).eq("id", document_id).execute()


def get_documents_by_student(student_id: str) -> list:
    res = supabase.table("documents").select("*").eq("student_id", student_id).order("created_at", desc=True).execute()
    return res.data


def get_documents_by_category(category_id: int) -> list:
    res = supabase.table("documents").select("*").eq("category_id", category_id).order("created_at", desc=True).execute()
    return res.data


def delete_document(document_id: int) -> None:
    res = supabase.table("documents").select("file_path").eq("id", document_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    file_path = res.data[0]["file_path"]

    # 청크 먼저 삭제 (FK 제약)
    try:
        supabase.table("document_chunks").delete().eq("document_id", document_id).execute()
    except Exception as e:
        logger.error("청크 삭제 실패: %s", e)
        raise HTTPException(status_code=500, detail="문서 청크 삭제 중 오류가 발생했습니다.")

    # 문서 레코드 삭제
    try:
        supabase.table("documents").delete().eq("id", document_id).execute()
    except Exception as e:
        logger.error("문서 DB 삭제 실패: %s", e)
        raise HTTPException(status_code=500, detail="문서 정보 삭제 중 오류가 발생했습니다.")

    # 스토리지 파일 삭제
    try:
        supabase.storage.from_("documents").remove([file_path])
    except Exception as e:
        logger.warning("스토리지 파일 삭제 실패 (이미 없을 수 있음): %s", e)
